

# region 创建word，设置格式和打开

'''
from docx import Document
from docx.shared import Cm, Pt

from docx.document import Document as Doc
import os

# 创建代表Word文档的Doc对象
document = Document()  # type: Doc
# 添加大标题
document.add_heading('快快乐乐学Python', 0)
# 添加段落
p = document.add_paragraph('Python是一门非常流行的编程语言，它')

run = p.add_run('简单')
run.bold = True
run.font.size = Pt(18)
p.add_run('而且')
run = p.add_run('优雅')
run.font.size = Pt(18)
run.underline = True
p.add_run('。')

# 添加一级标题
document.add_heading('Heading, level 1', level=1)
# 添加带样式的段落
document.add_paragraph('Intense quote', style='Intense Quote')
# 添加无序列表
document.add_paragraph(
    'first item in unordered list', style='List Bullet'
)
document.add_paragraph(
    'second item in ordered list', style='List Bullet'
)
# 添加有序列表
document.add_paragraph(
    'first item in ordered list', style='List Number'
)
document.add_paragraph(
    'second item in ordered list', style='List Number'
)

# 添加图片（注意路径和图片必须要存在）
document.add_picture('resources/guido.jpg', width=Cm(5.2))

# 添加分节符
document.add_section()

records = (
    ('骆昊', '男', '1995-5-5'),
    ('孙美丽', '女', '1992-2-2')
)
# 添加表格
table = document.add_table(rows=1, cols=3)
table.style = 'Dark List'
hdr_cells = table.rows[0].cells
hdr_cells[0].text = '姓名'
hdr_cells[1].text = '性别'
hdr_cells[2].text = '出生日期'
# 为表格添加行
for name, sex, birthday in records:
    row_cells = table.add_row().cells
    row_cells[0].text = name
    row_cells[1].text = sex
    row_cells[2].text = birthday


# 添加分页符
document.add_page_break()


# 保存文档
document.save('demo.docx')


# 自动打开文档
try:
    if os.name == 'nt':  # Windows 系统
        os.startfile(r'D:\VSCODE\python100exe\resources\demo.docx')
    elif os.name == 'posix':  # macOS 和 Linux 系统
        import subprocess
        subprocess.call(('open', file_path))
    else:
        print("不支持的操作系统")
except Exception as e:
    print(f"打开文件时出错: {e}")
'''
# endregion



# region  读取word内容

'''
from docx import Document
from docx.document import Document as Doc

doc = Document('resources/枚举enum.docx')  # type: Doc
for no, p in enumerate(doc.paragraphs):
    print(no, p.text)
'''
#endregion 


# region 从离职模板里读取，插入数据后生成个人的离职证明

'''
from docx import Document
from docx.document import Document as Doc

# 将真实信息用字典的方式保存在列表中
employees = [
    {
        'name': '骆昊',
        'id': '100200198011280001',
        'sdate': '2008年3月1日',
        'edate': '2012年2月29日',
        'department': '产品研发',
        'position': '架构师',
        'company': '成都华为技术有限公司'
    },
    {
        'name': '王大锤',
        'id': '510210199012125566',
        'sdate': '2019年1月1日',
        'edate': '2021年4月30日',
        'department': '产品研发',
        'position': 'Python开发工程师',
        'company': '成都谷道科技有限公司'
    },
    {
        'name': '李元芳',
        'id': '2102101995103221599',
        'sdate': '2020年5月10日',
        'edate': '2021年3月5日',
        'department': '产品研发',
        'position': 'Java开发工程师',
        'company': '同城企业管理集团有限公司'
    },
]
# 对列表进行循环遍历，批量生成Word文档 
for emp_dict in employees:
    # 读取离职证明模板文件
    doc = Document('resources/离职证明模板.docx')  # type: Doc
    # 循环遍历所有段落寻找占位符
    for p in doc.paragraphs:
        if '{' not in p.text:
            continue
        # 不能直接修改段落内容，否则会丢失样式
        # 所以需要对段落中的元素进行遍历并进行查找替换
        for run in p.runs:
            if '{' not in run.text:
                continue
            # 将占位符换成实际内容
            start, end = run.text.find('{'), run.text.find('}')
            key, place_holder = run.text[start + 1:end], run.text[start:end + 1]
            run.text = run.text.replace(place_holder, emp_dict[key])
    # 每个人对应保存一个Word文档
    doc.save(f'resources/{emp_dict["name"]}离职证明.docx')
'''
#endregion 



# region 创建设置ppt格式，打开ppt

from pptx import Presentation
import os

# 创建幻灯片对象
pres = Presentation()

# 选择母版添加一页
title_slide_layout = pres.slide_layouts[0]
slide = pres.slides.add_slide(title_slide_layout)

# 获取页面上所有形状
shapes = slide.shapes
# 获取标题和主体
title_shape = shapes.title
body_shape = shapes.placeholders[1]
# 编辑标题
title_shape.text = 'Introduction'

# 编辑主体内容
tf = body_shape.text_frame
tf.text = 'History of Python'
# 添加一个一级段落
p = tf.add_paragraph()
p.text = 'X\'max 1989'
p.level = 1
# 添加一个二级段落
p = tf.add_paragraph()
p.text = 'Guido began to write interpreter for Python.'
p.level = 2

# 保存幻灯片
pres.save('resources/test.pptx')
print(f"文件已保存到：resources")

# 自动打开文档
try:
    if os.name == 'nt':  # Windows 系统
        os.startfile(r'D:\VSCODE\python100exe\resources\test.pptx')
    elif os.name == 'posix':  # macOS 和 Linux 系统
        import subprocess
        subprocess.call(('open', resources/test.pptx))
    else:
        print("不支持的操作系统")
except Exception as e:
    print(f"打开文件时出错: {e}")

#endregion 